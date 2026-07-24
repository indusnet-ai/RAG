import io, base64, hashlib, logging, re
from dataclasses import dataclass
from datetime import datetime
from typing import List, Dict, Any, Optional
from pathlib import Path
from pdf2image import convert_from_bytes
from PIL import Image
from pptx import Presentation
import pptx.util
import pymupdf
from openai import OpenAI
from dotenv import load_dotenv
import os
import pandas as pd
from docx import Document
from concurrent.futures import ThreadPoolExecutor, as_completed
import time
import openai
from langsmith import wrappers




# Load variables from .env into environment
load_dotenv()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
    logger.info("✅ HEIC/HEIF support enabled")
except ImportError:
    logger.warning("⚠️ pillow-heif not installed - HEIC/HEIF files won't work")
client = OpenAI()

# Try to import FastChunker and TokenChunker
try:
    from chonkie import FastChunker, TokenChunker
    FAST_CHUNKER_AVAILABLE = True
    TOKEN_CHUNKER_AVAILABLE = True
    logger.info("✅ FastChunker and TokenChunker available from chonkie")
except ImportError:
    try:
        from chonkie import TokenChunker
        TOKEN_CHUNKER_AVAILABLE = True
        FAST_CHUNKER_AVAILABLE = False
        FastChunker = None
        logger.info("✅ TokenChunker available from chonkie")
    except ImportError:
        TOKEN_CHUNKER_AVAILABLE = False
        FAST_CHUNKER_AVAILABLE = False
        FastChunker = None
        logger.info("⚠️  Chonkie chunkers not available")


# -------------------------------
# DocumentChunk
# -------------------------------
@dataclass
class DocumentChunk:
    """Represents a processed document chunk with metadata for citations"""
    content: str
    source_file: str
    source_type: str  # 'pdf', 'txt', 'web', 'audio', 'ocr_pdf', 'docx'
    page_number: Optional[int] = None
    chunk_index: int = 0
    start_char: Optional[int] = None
    end_char: Optional[int] = None
    metadata: Dict[str, Any] = None
    chunk_id: str = ""

    def __post_init__(self):
        if not self.chunk_id:
            self.chunk_id = self._generate_chunk_id()
        if self.metadata is None:
            self.metadata = {}

    def _generate_chunk_id(self) -> str:
        content_hash = hashlib.md5(self.content.encode()).hexdigest()[:8]
        return f"{self.source_type}_{self.chunk_index}_{content_hash}"

    def get_citation_info(self) -> Dict[str, Any]:
        citation = {
            'source': self.source_file,
            'type': self.source_type,
            'chunk_id': self.chunk_id,
            'chunk_index': self.chunk_index
        }
        if self.page_number:
            citation['page'] = self.page_number
        if self.start_char or self.end_char:
            citation['char_range'] = f"{self.start_char}-{self.end_char}"
        citation.update(self.metadata)
        return citation


class DocumentProcessor:
    def __init__(
        self, 
        chunk_size: int = 1500, 
        chunk_overlap: int = 400, 
        max_workers: int = 5,
        pages_per_batch: int = 3,  # Pages to batch together for OCR
        use_fast_chunker: bool = True,  # Use memchunk for 3000x faster chunking
        use_token_chunker: bool = True  # Use TokenChunker for token-based chunking
    ):
        """
        Initialize DocumentProcessor with batched parallel OCR support, token chunking and optional fast chunking
        """
        self.max_workers = max_workers
        self.pages_per_batch = pages_per_batch
        self.supported_formats = {
            # Documents
            '.pdf', '.txt', '.md', '.docx', '.pptx',
            # Images
            '.png', '.jpg', '.jpeg', '.bmp', '.gif', '.webp', 
            '.tif', '.tiff', '.heic', '.heif', '.avif', '.ico', '.jp2',
            # audio
            '.mp3','.webm','wav',
            # Spreadsheets
            '.csv', '.xlsx', '.xls', '.xlsm'}       
        
        # Check if we should use TokenChunker
        self.use_token_chunker = use_token_chunker and TOKEN_CHUNKER_AVAILABLE
        self.token_chunker = None
        
        if self.use_token_chunker:
            # Set default token sizes if defaulted to character limits
            self.chunk_size = 400 if chunk_size == 1500 else chunk_size
            self.chunk_overlap = 80 if chunk_overlap == 400 else chunk_overlap
            try:
                self.token_chunker = TokenChunker(
                    tokenizer='gpt2',
                    chunk_size=self.chunk_size,
                    chunk_overlap=self.chunk_overlap
                )
                logger.info(f"✅ TokenChunker (GPT-2) initialized: size={self.chunk_size} tokens, overlap={self.chunk_overlap} tokens")
            except Exception as e:
                logger.warning(f"⚠️ Failed to initialize TokenChunker: {e}. Falling back to standard/fast chunking.")
                self.use_token_chunker = False
                self.chunk_size = chunk_size
                self.chunk_overlap = chunk_overlap
        else:
            self.chunk_size = chunk_size
            self.chunk_overlap = chunk_overlap

        # Initialize FastChunker if requested and available (fallback/alternative)
        self.use_fast_chunker = use_fast_chunker and FAST_CHUNKER_AVAILABLE and not self.use_token_chunker
        self.fast_chunker = None
        
        if self.use_fast_chunker:
            try:
                self.fast_chunker = FastChunker(
                    chunk_size=self.chunk_size
                )
                logger.info(f"✅ FastChunker (memchunk) initialized - 3000x faster chunking enabled!")
            except Exception as e:
                logger.warning(f"⚠️ Failed to initialize FastChunker: {e}. Falling back to standard chunking.")
                self.use_fast_chunker = False
        
        logger.info(f"🚀 DocumentProcessor initialized:")
        logger.info(f"   - Chunk size: {self.chunk_size}")
        logger.info(f"   - Chunk overlap: {self.chunk_overlap}")
        logger.info(f"   - Parallel workers: {self.max_workers}")
        logger.info(f"   - Pages per batch: {self.pages_per_batch}")
        logger.info(f"   - Token chunking: {'✅ ENABLED' if self.use_token_chunker else '⚠️ DISABLED'}")
        logger.info(f"   - Fast chunking: {'✅ ENABLED (memchunk)' if self.use_fast_chunker else '⚠️ DISABLED'}")

    def _is_page_image_based(self, page) -> bool:
        """Decide if an individual PDF page is image-based."""
        try:
            text = page.get_text().strip()
            return len(text) < 50  # threshold; tweak as needed
        except:
            return True

    def _ocr_batch_pages(
        self, 
        page_images: List[tuple], 
        batch_id: int
    ) -> Dict[str, Any]:
        """
        🆕 Perform OCR on MULTIPLE pages in a SINGLE API call
        
        Args:
            page_images: List of (page_image, page_num) tuples
            batch_id: Batch identifier for logging
            
        Returns:
            Dict with page_num → text mapping
        """
        try:
            start_time = time.time()
            
            # Build content with multiple images
            content = [
                {
                    "type": "input_text",
                    "text": (
                        "Extract all text from each page exactly as written. "
                        "For each page, start with '--- PAGE X ---' where X is the page number. "
                        "Do not translate. Preserve original formatting and language."
                    )
                }
            ]
            
            # Add all images to single request
            for page_image, page_num in page_images:
                buffered = io.BytesIO()
                page_image.save(buffered, format="PNG")
                img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
                
                content.append({
                    "type": "input_text",
                    "text": f"\n--- PAGE {page_num} ---"
                })
                content.append({
                    "type": "input_image",
                    "image_url": f"data:image/png;base64,{img_b64}"
                })
            
            # Single API call for all pages in batch
            response = client.responses.create(
                model="gpt-4.1-mini",
                input=[{"role": "user", "content": content}],
            )
            
            # Parse response to extract text per page
            full_text = response.output_text
            page_results = {}
            
            # Split by page markers
            page_sections = re.split(r'---\s*PAGE\s+(\d+)\s*---', full_text)
            
            # Extract text for each page
            for i in range(1, len(page_sections), 2):
                if i + 1 < len(page_sections):
                    page_num = int(page_sections[i])
                    page_text = page_sections[i + 1].strip()
                    page_results[page_num] = page_text
            
            # Fallback: if parsing failed, assign all text to first page
            if not page_results:
                logger.warning(f"⚠️ Batch {batch_id}: Failed to split pages, using full text")
                for _, page_num in page_images:
                    page_results[page_num] = full_text
            
            elapsed = time.time() - start_time
            logger.info(
                f"✅ Batch {batch_id} ({len(page_images)} pages) "
                f"completed in {elapsed:.2f}s"
            )
            
            return {
                "batch_id": batch_id,
                "pages": page_results,
                "success": True,
                "elapsed_time": elapsed,
                "page_count": len(page_images)
            }
            
        except Exception as e:
            logger.error(f"❌ Error in batch {batch_id}: {str(e)}")
            # Return error result for all pages in this batch
            error_results = {
                page_num: f"[OCR Error on page {page_num}: {str(e)}]"
                for _, page_num in page_images
            }
            return {
                "batch_id": batch_id,
                "pages": error_results,
                "success": False,
                "error": str(e),
                "page_count": len(page_images)
            }

    def _ocr_single_page(
        self, 
        page_image: Image.Image, 
        page_num: int
    ) -> Dict[str, Any]:
        """
        Legacy method: Perform OCR on a single page
        (Kept for compatibility with single-image processing)
        """
        try:
            start_time = time.time()
            
            buffered = io.BytesIO()
            page_image.save(buffered, format="PNG")
            img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
            
            response = client.responses.create(
                model="gpt-4.1-mini",
                input=[
                    {
                        "role": "user",
                        "content": [
                             {
                                "type": "input_text",
                                "text": """Perform the following tasks conditionally:\n\n\
                                    1. OCR (ONLY IF text is present): Extract all visible text exactly as written. Preserve original formatting and language.\n\n\
                                    2. Visual Analysis (ONLY IF non-textual visual elements are present): Describe meaningful non-textual elements such as objects, charts, diagrams, spatial layout, or relationships. Ignore purely decorative elements.\n\n\
                                        ONLY RETURN THE OUTPUT""" },
                            {
                                "type": "input_image",
                                "image_url": f"data:image/png;base64,{img_b64}"
                            }
                        ],
                    }
                ],
            )
            
            page_text = response.output_text
            elapsed = time.time() - start_time
            logger.info(f"✅ Page {page_num} OCR completed in {elapsed:.2f}s")
            
            return {
                "page_num": page_num,
                "text": page_text,
                "success": True,
                "elapsed_time": elapsed
            }
        except Exception as e:
            logger.error(f"❌ Error OCR'ing page {page_num}: {str(e)}")
            return {
                "page_num": page_num,
                "text": f"[OCR Error on page {page_num}: {str(e)}]",
                "success": False,
                "error": str(e)
            }

    def process_document(self, file_path: str) -> List[DocumentChunk]:
        """Main entry point for document processing"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if file_path.suffix.lower() not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_path.suffix}")
        
        logger.info(f"📄 Processing: {file_path.name}")
        
        try:
            # ADD THIS before the existing PDF check:
            if file_path.suffix.lower() in {'.xlsx', '.xls', '.xlsm', '.csv'}:
                return self._process_excel(file_path)
            elif file_path.suffix.lower() == '.pdf':
                return self._process_pdf_hybrid_batched(file_path)
            elif file_path.suffix.lower() == '.docx':
                logger.info(f"📘 DOCX → Using python-docx")
                return self._process_docx(file_path)
            elif file_path.suffix.lower() in {'.pptx', '.ppt'}:  # ← ADD THIS BLOCK
                logger.info(f"📊 PPTX → Using python-pptx")
                return self._process_pptx(file_path)
            elif file_path.suffix.lower() in {'.png', '.jpg', '.jpeg', '.bmp', '.gif', '.webp', 
            '.tif', '.tiff', '.heic', '.heif', '.avif', '.ico', '.jp2'}:
                logger.info(f"🖼️ Image → Using OCR")
                return self._process_ocr_image(file_path)
            elif file_path.suffix.lower() in {'.txt', '.md'}:
                logger.info(f"📝 Text file → Direct reading")
                return self._process_text_file(file_path)
        except Exception as e:
            logger.error(f"❌ Error processing {file_path.name}: {str(e)}")
            raise

    def _process_pdf_hybrid_batched(self, file_path: Path) -> List[DocumentChunk]:
        """
        🆕 BATCHED Hybrid PDF processing:
        - Extract text pages immediately
        - Batch image pages together (e.g., 3-5 per API call)
        - Process batches in parallel
        - Maintain correct page order
        """
        doc = pymupdf.open(file_path)
        total_pages = len(doc)
        logger.info(f"📄 Batched Hybrid PDF → {total_pages} pages")
        
        text_pages = {}  # page_num → extracted text
        ocr_page_images = {}  # page_num → PIL image
        
        # --------------------------------
        # 1️⃣ DETECT PAGE TYPE + PREPROCESS
        # --------------------------------
        for page_idx in range(total_pages):
            page = doc.load_page(page_idx)
            txt = page.get_text().strip()
            
            # Check if text is just boilerplate/footer
            is_boilerplate = any(w in txt for w in ["Follow Chandra", "Repost to help", "GET HIRED", "THAT'S A WRAP"])
            if len(txt) >= 150 and not is_boilerplate:
                # TEXT PAGE
                text_pages[page_idx + 1] = {
                    "text": txt,
                    "method": "pymupdf"
                }
            else:
                # IMAGE PAGE → prepare for OCR
                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                ocr_page_images[page_idx + 1] = img
        
        doc.close()
        
        # --------------------------------
        # 2️⃣ BATCH OCR PAGES
        # --------------------------------
        ocr_results = {}
        
        if ocr_page_images:
            logger.info(
                f"🖼️ {len(ocr_page_images)} pages require OCR → "
                f"batching {self.pages_per_batch} pages per call..."
            )
            
            # Create batches
            page_items = list(ocr_page_images.items())
            batches = []
            
            for i in range(0, len(page_items), self.pages_per_batch):
                batch = [
                    (img, page_num) 
                    for page_num, img in page_items[i:i + self.pages_per_batch]
                ]
                batches.append(batch)
            
            logger.info(f"📦 Created {len(batches)} batches")
            
            # Process batches in parallel
            with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                futures = {
                    executor.submit(self._ocr_batch_pages, batch, batch_id): batch_id
                    for batch_id, batch in enumerate(batches)
                }
                
                completed = 0
                total_batches = len(batches)
                
                for future in as_completed(futures):
                    batch_id = futures[future]
                    try:
                        result = future.result()
                        ocr_results.update(result["pages"])
                        completed += 1
                        logger.info(
                            f"📊 Progress: {completed}/{total_batches} batches completed"
                        )
                    except Exception as e:
                        logger.error(f"❌ Exception in batch {batch_id}: {str(e)}")
                        completed += 1
        
        # --------------------------------
        # 3️⃣ MERGE RESULTS PAGE-WISE
        # --------------------------------
        combined_pages = {}
        for p in range(1, total_pages + 1):
            if p in text_pages:
                combined_pages[p] = text_pages[p]
            elif p in ocr_results:
                combined_pages[p] = {
                    "text": ocr_results[p],
                    "method": "ocr_batched"
                }
            else:
                combined_pages[p] = {
                    "text": "",
                    "method": "missing"
                }
        
        # --------------------------------
        # 4️⃣ CREATE CHUNKS PAGE-WISE
        # --------------------------------
        chunks = []
        for page_num, info in combined_pages.items():
            page_text = info["text"]
            method_used = info["method"]
            
            section = "Introduction"
            if page_text:
                lines = [line.strip() for line in page_text.split('\n') if line.strip()]
                for line in lines[:3]:
                    if len(line) < 100 and (
                        line.isupper() or 
                        re.match(r'^\d+(\.\d+)*\s+[a-zA-Z]', line) or 
                        line.startswith('#')
                    ):
                        section = line.lstrip('#').strip()
                        break

            metadata = {
                "processing_method": method_used,
                "processed_at": datetime.now().isoformat(),
                "page": page_num,
                "page_number": page_num,
                "section": section,
                "title": file_path.name,
                "ocr_engine": "openai_vision" if "ocr" in method_used else None,
                "batched": "ocr" in method_used,
                "pages_per_batch": self.pages_per_batch if "ocr" in method_used else None
            }
            
            page_chunks = self._create_chunks_from_text(
                text=page_text,
                source_file=file_path.name,
                source_type="pdf_hybrid_batched",
                page_number=page_num,
                additional_metadata=metadata
            )
            chunks.extend(page_chunks)
        
        logger.info(f"✅ Batched Hybrid PDF complete: {len(chunks)} chunks generated")
        return chunks

    def _process_docx(self, file_path: Path) -> List[DocumentChunk]:
        """Process Microsoft Word .docx files"""
        try:
            doc = Document(file_path)
            full_text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
            
            content = '\n\n'.join(full_text)
            
            metadata = {
                'file_size': file_path.stat().st_size,
                'paragraph_count': len(doc.paragraphs),
                'processed_at': datetime.now().isoformat(),
                'processing_method': 'python-docx'
            }
            
            try:
                core_props = doc.core_properties
                if core_props.author:
                    metadata['author'] = core_props.author
                if core_props.title:
                    metadata['title'] = core_props.title
                if core_props.created:
                    metadata['created'] = core_props.created.isoformat()
                if core_props.modified:
                    metadata['modified'] = core_props.modified.isoformat()
            except:
                pass
            
            chunks = self._create_chunks_from_text(
                content,
                file_path.name,
                source_type='docx',
                page_number=None,
                additional_metadata=metadata
            )
            
            logger.info(f"✅ DOCX: {len(chunks)} chunks from {len(doc.paragraphs)} paragraphs")
            return chunks
        except Exception as e:
            logger.error(f"❌ Error processing DOCX: {str(e)}")
            raise
    
    def _process_pptx(self, file_path: Path) -> List[DocumentChunk]:
        """
        🆕 Hybrid PPTX Processing with OCR Support
        
        Intelligently processes PowerPoint slides by:
        1. Extracting text from text boxes and shapes
        2. Checking if slide has embedded images
        3. Using OCR on slides with images or low text content
        4. Combining results for complete content extraction
        
        This matches the PDF hybrid processing approach!
        """
        try:
            
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            
            logger.info(f"📊 Processing PPTX (hybrid mode): {total_slides} slides")
            
            # ============================================
            # STEP 1: Analyze each slide
            # ============================================
            
            slides_to_process = []
            slides_to_ocr = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                # Extract text from shapes
                slide_text = self._extract_text_from_slide(slide)
                
                # Check if slide has images
                has_images = self._slide_has_images(slide)
                
                # Check if slide has enough text
                text_length = len(slide_text.strip())
                
                # Decision logic (similar to PDF)
                if text_length < 50 or has_images:
                    # Slide needs OCR (low text or has images)
                    slides_to_ocr.append((slide, slide_num))
                    logger.info(f"  Slide {slide_num}: Needs OCR (text={text_length} chars, images={has_images})")
                else:
                    # Text extraction is sufficient
                    slides_to_process.append({
                        'slide_num': slide_num,
                        'text': slide_text,
                        'method': 'text_extraction'
                    })
                    logger.info(f"  Slide {slide_num}: Text extraction (text={text_length} chars)")
            
            logger.info(f"📊 Strategy: {len(slides_to_process)} text-based, {len(slides_to_ocr)} need OCR")
            
            # ============================================
            # STEP 2: Render slides that need OCR
            # ============================================
            
            ocr_results = {}
            if slides_to_ocr:
                logger.info(f"🖼️  Rendering {len(slides_to_ocr)} slides for OCR...")
                
                # Render slides as images
                slide_images = []
                for slide, slide_num in slides_to_ocr:
                    try:
                        # Render slide to image
                        img = self._render_slide_to_image(file_path, slide_num)
                        if img:
                            slide_images.append((img, slide_num))
                    except Exception as e:
                        logger.error(f"Failed to render slide {slide_num}: {e}")
                
                # ============================================
                # STEP 3: Batch OCR on rendered slides
                # ============================================
                
                if slide_images:
                    logger.info(f"🔍 Running OCR on {len(slide_images)} slides...")
                    
                    # Batch slides for OCR (like PDF does)
                    batches = [
                        slide_images[i:i + self.pages_per_batch]
                        for i in range(0, len(slide_images), self.pages_per_batch)
                    ]
                    
                    with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                        futures = {
                            executor.submit(self._ocr_batch_pages, batch, batch_id): batch_id
                            for batch_id, batch in enumerate(batches)
                        }
                        
                        for future in as_completed(futures):
                            batch_id = futures[future]
                            try:
                                batch_results = future.result()
                                ocr_results.update(batch_results)
                            except Exception as e:
                                logger.error(f"Batch {batch_id} failed: {e}")
            
            # ============================================
            # STEP 4: Combine text and OCR results
            # ============================================
            
            combined_slides = {}
            
            # Add text-extracted slides
            for slide_info in slides_to_process:
                combined_slides[slide_info['slide_num']] = {
                    'text': slide_info['text'],
                    'method': 'text_extraction'
                }
            
            # Add OCR slides
            for slide, slide_num in slides_to_ocr:
                # Get both extracted text and OCR text
                extracted_text = self._extract_text_from_slide(slide)
                ocr_text = ocr_results.get(slide_num, "")
                
                # Combine (OCR text usually more complete for image-heavy slides)
                if extracted_text and ocr_text:
                    # Both available - combine intelligently
                    combined_text = f"{extracted_text}\n\n[Image Content]\n{ocr_text}"
                    method = "hybrid_text_ocr"
                elif ocr_text:
                    combined_text = ocr_text
                    method = "ocr_only"
                else:
                    combined_text = extracted_text
                    method = "text_fallback"
                
                combined_slides[slide_num] = {
                    'text': combined_text,
                    'method': method
                }
            
            # ============================================
            # STEP 5: Create chunks from all slides
            # ============================================
            
            chunks = []
            for slide_num in sorted(combined_slides.keys()):
                slide_info = combined_slides[slide_num]
                
                metadata = {
                    'slide_number': slide_num,
                    'processing_method': slide_info['method'],
                    'processed_at': datetime.now().isoformat(),
                    'ocr_engine': 'openai_vision' if 'ocr' in slide_info['method'] else None,
                    'file_size': file_path.stat().st_size,
                    'total_slides': total_slides
                }
                
                slide_chunks = self._create_chunks_from_text(
                    text=slide_info['text'],
                    source_file=file_path.name,
                    source_type='pptx_hybrid',
                    page_number=slide_num,
                    additional_metadata=metadata
                )
                
                chunks.extend(slide_chunks)
            
            logger.info(f"✅ Hybrid PPTX complete: {len(chunks)} chunks from {total_slides} slides")
            return chunks
            
        except Exception as e:
            logger.error(f"❌ Error in hybrid PPTX processing: {e}")
            logger.info("⚠️  Falling back to standard text-only processing...")
            return self._process_pptx(file_path)  # Fallback to old method
    
    def _extract_text_from_slide(self, slide) -> str:
        """
        Extract all text from a slide's shapes and tables
        """
        slide_text = []
        
        for shape in slide.shapes:
            # Regular text
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            
            # Tables
            if shape.has_table:
                table_text = []
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(' | '.join(row_text))
                if table_text:
                    slide_text.append('\n'.join(table_text))
        
        return '\n\n'.join(slide_text)
    
    def _slide_has_images(self, slide) -> bool:
        """
        Check if slide contains images
        """
        try:
            for shape in slide.shapes:
                # Check for picture shapes
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    return True
                
                # Check for grouped shapes containing images
                if hasattr(shape, 'shapes'):  # Group shape
                    for subshape in shape.shapes:
                        if subshape.shape_type == 13:
                            return True
            
            return False
        except:
            return False
    
    def _render_slide_to_image(self, pptx_path: Path, slide_num: int) -> Optional[Image.Image]:
        """
        Render a PPTX slide to an image
        
        Options:
        1. Use LibreOffice (if installed)
        2. Use python-pptx-image (if available)
        3. Use Aspose.Slides (commercial)
        
        For now, using LibreOffice approach (most common)
        """
        try:
            import subprocess
            import tempfile
            
            # Create temp directory for output
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_dir_path = Path(temp_dir)
                
                # Convert PPTX to images using LibreOffice
                # Note: This requires LibreOffice to be installed
                cmd = [
                    'soffice',
                    '--headless',
                    '--convert-to', 'png',
                    '--outdir', str(temp_dir_path),
                    str(pptx_path)
                ]
                
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    timeout=30
                )
                
                if result.returncode == 0:
                    # Find the generated PNG for this slide
                    png_files = list(temp_dir_path.glob('*.png'))
                    if png_files and slide_num <= len(png_files):
                        # Slides are numbered from 1, files from 0
                        img_path = sorted(png_files)[slide_num - 1]
                        return Image.open(img_path)
                
                return None
                
        except Exception as e:
            logger.debug(f"Could not render slide {slide_num}: {e}")
            return None
    
    def _process_excel(self, file_path: Path) -> List[DocumentChunk]:
        """Process Excel files (.xlsx, .xls, .xlsm, .csv) with robust error handling"""
        import pandas as pd
        
        try:
            logger.info(f"📊 Processing Excel file: {file_path.name}")
            
            file_ext = file_path.suffix.lower()
            
            if file_ext == '.csv':
                # ✅ ROBUST CSV PARSING WITH MULTIPLE STRATEGIES
                encodings_to_try = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252', 'utf-16']
                df = None
                successful_encoding = None
                parsing_strategy = None
                
                # STRATEGY 1: Standard parsing
                for encoding in encodings_to_try:
                    try:
                        df = pd.read_csv(file_path, encoding=encoding)
                        successful_encoding = encoding
                        parsing_strategy = 'standard'
                        logger.info(f"✅ CSV loaded with {encoding} (standard parsing)")
                        break
                    except (UnicodeDecodeError, UnicodeError):
                        continue
                    except pd.errors.ParserError:
                        continue
                
                # STRATEGY 2: Skip bad lines
                if df is None:
                    logger.warning("⚠️  Standard parsing failed, trying with on_bad_lines='skip'")
                    for encoding in encodings_to_try:
                        try:
                            df = pd.read_csv(
                                file_path, 
                                encoding=encoding,
                                on_bad_lines='skip'  # ✅ Skip malformed rows
                            )
                            successful_encoding = encoding
                            parsing_strategy = 'skip_bad_lines'
                            logger.info(f"✅ CSV loaded with {encoding} (skipped bad lines)")
                            break
                        except (UnicodeDecodeError, UnicodeError):
                            continue
                        except Exception as e:
                            logger.debug(f"Encoding {encoding} failed: {e}")
                            continue
                
                # STRATEGY 3: Use Python engine (more flexible, slower)
                if df is None:
                    logger.warning("⚠️  Trying Python engine with flexible parsing")
                    for encoding in encodings_to_try:
                        try:
                            df = pd.read_csv(
                                file_path,
                                encoding=encoding,
                                engine='python',  # ✅ More forgiving parser
                                on_bad_lines='skip',
                                sep=None  # ✅ Auto-detect delimiter
                            )
                            successful_encoding = encoding
                            parsing_strategy = 'python_engine'
                            logger.info(f"✅ CSV loaded with {encoding} (Python engine)")
                            break
                        except Exception as e:
                            logger.debug(f"Python engine with {encoding} failed: {e}")
                            continue
                
                # STRATEGY 4: Read as plain text (last resort)
                if df is None:
                    logger.warning("⚠️  All CSV parsing failed, reading as plain text")
                    try:
                        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                            text_content = f.read()
                        
                        # Return as single text chunk instead of dataframe
                        chunks = self._create_chunks_from_text(
                            text=text_content,
                            source_file=file_path.name,
                            source_type='csv_text',
                            page_number=None,
                            additional_metadata={
                                'parsing_method': 'plain_text_fallback',
                                'warning': 'CSV structure was malformed, processed as plain text'
                            }
                        )
                        logger.info(f"✅ CSV processed as plain text: {len(chunks)} chunks")
                        return chunks
                        
                    except Exception as e:
                        raise ValueError(
                            f"Failed to read CSV with any method. File may be corrupted. Error: {str(e)}"
                        )
                
                sheets_data = {'Sheet1': df}
                
            else:
                # Excel files (.xlsx, .xls, .xlsm) - standard processing
                with pd.ExcelFile(file_path) as excel_file:
                    sheets_data = {  
                        sheet: pd.read_excel(excel_file, sheet_name=sheet) 
                        for sheet in excel_file.sheet_names
                    }   
                successful_encoding = 'excel_binary'
                parsing_strategy = 'excel'
            
            all_chunks = []
            chunk_index = 0
            
            for sheet_name, df in sheets_data.items():
                logger.info(f"  Processing sheet: {sheet_name} ({len(df)} rows × {len(df.columns)} cols)")
                
                sheet_text = self._excel_to_text(df, sheet_name)
                
                if not sheet_text.strip():
                    continue
                
                chunks = self._create_chunks_from_text(
                    text=sheet_text,
                    source_file=file_path.name,
                    source_type='excel',
                    page_number=None,
                    additional_metadata={
                        'sheet_name': sheet_name,
                        'rows': len(df),
                        'columns': len(df.columns),
                        'encoding': successful_encoding,
                        'parsing_strategy': parsing_strategy,  # ✅ Track how it was parsed
                        'file_extension': file_ext
                    }
                )
                
                for chunk in chunks:
                    chunk.chunk_index = chunk_index
                    chunk_index += 1
                
                all_chunks.extend(chunks)
            
            logger.info(f"✅ Excel processing complete: {len(all_chunks)} total chunks")
            return all_chunks
            
        except Exception as e:
            logger.error(f"❌ Error processing Excel: {e}", exc_info=True)
            raise

    def _excel_to_text(self, df, sheet_name: str) -> str:
        """Convert Excel DataFrame to formatted text"""        
        lines = [f"## Sheet: {sheet_name}", ""]
        
        # Add headers
        lines.append("| " + " | ".join(str(col) for col in df.columns) + " |")
        lines.append("|" + "|".join(["---" for _ in df.columns]) + "|")
        
        # Add rows (limit to 1000)
        for idx, row in df.head(1000).iterrows():
            cells = []
            for val in row:
                if pd.isna(val):
                    cells.append("")
                else:
                    cell_str = str(val).strip()
                    if len(cell_str) > 200:
                        cell_str = cell_str[:197] + "..."
                    cells.append(cell_str)
            lines.append("| " + " | ".join(cells) + " |")
        
        if len(df) > 1000:
            lines.append(f"\n*[Showing first 1000 of {len(df)} rows]*")
        
        return "\n".join(lines)

    def _process_ocr_image(self, file_path: Path) -> List[DocumentChunk]:
        """Process image files using OCR (single image - no batching)"""
        try:
            with open(file_path, 'rb') as file:
                file_content = file.read()
            
            image = Image.open(io.BytesIO(file_content))
            result = self._ocr_single_page(image, 1)
            ocr_text = result["text"]
            
            metadata = {
                'file_size': file_path.stat().st_size,
                'processed_at': datetime.now().isoformat(),
                'processing_method': 'ocr',
                'ocr_engine': 'openai_vision'
            }
            
            chunks = self._create_chunks_from_text(
                ocr_text,
                file_path.name,
                source_type='ocr_image',
                page_number=1,
                additional_metadata=metadata
            )
            
            logger.info(f"✅ Image OCR: {len(chunks)} chunks")
            return chunks
        except Exception as e:
            logger.error(f"❌ Error processing image: {str(e)}")
            raise

    def _process_text_file(self, file_path: Path) -> List[DocumentChunk]:
        """Process regular text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            metadata = {
                'file_size': file_path.stat().st_size,
                'encoding': 'utf-8',
                'processed_at': datetime.now().isoformat(),
                'processing_method': 'direct'
            }
            
            chunks = self._create_chunks_from_text(
                content,
                file_path.name,
                source_type='txt',
                page_number=None,
                additional_metadata=metadata
            )
            
            logger.info(f"✅ Text file: {len(chunks)} chunks")
            return chunks
        except Exception as e:
            logger.error(f"❌ Error processing text file: {str(e)}")
            raise
        
    def _apply_overlap(self, chunks: List[DocumentChunk]) -> List[DocumentChunk]:
        if self.chunk_overlap <= 0:
            return chunks

        overlapped = []
        for i, chunk in enumerate(chunks):
            if i == 0:
                overlapped.append(chunk)
                continue

            prev = chunks[i - 1]
            overlap_text = prev.content[-self.chunk_overlap:]

            chunk.content = overlap_text + chunk.content
            chunk.start_char = max(
                0, (chunk.start_char or 0) - self.chunk_overlap
            )

            overlapped.append(chunk)

        return overlapped
    

    def _sanitize_text(self, text: str) -> str:
        """
        Remove characters that can't be stored in PostgreSQL TEXT columns
        
        Args:
            text: Raw text that may contain problematic characters
            
        Returns:
            Cleaned text safe for database storage
        """
        # Remove NUL bytes
        text = text.replace('\x00', '')
        
        # Optionally remove other control characters (except newline, tab, carriage return)
        # import re
        # text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
        
        return text


    def _create_chunks_from_text(
        self,
        text: str,
        source_file: str,
        source_type: str,
        page_number: Optional[int] = None,
        additional_metadata: Dict[str, Any] = None
    ) -> List[DocumentChunk]:
        """
        Create standardized chunks from text
        
        Routes to TokenChunker (GPT-2 tokens), FastChunker (memchunk), or standard chunking
        """
        if not text.strip():
            return []
        
        text = self._sanitize_text(text)
        
        # Route to token, fast, or standard chunking
        if self.use_token_chunker and self.token_chunker:
            return self._create_chunks_token(
                text, source_file, source_type,
                page_number, additional_metadata
            )
        elif self.use_fast_chunker and self.fast_chunker:
            return self._create_chunks_fast(
                text, source_file, source_type, 
                page_number, additional_metadata
            )
        else:
            return self._create_chunks_standard(
                text, source_file, source_type,
                page_number, additional_metadata
            )

    def _create_chunks_token(
        self,
        text: str,
        source_file: str,
        source_type: str,
        page_number: Optional[int] = None,
        additional_metadata: Dict[str, Any] = None
    ) -> List[DocumentChunk]:
        """
        Token-based chunking using chonkie TokenChunker + post-processing overlap.
        Enhanced with semantic chunking by Heading, Section, and Architecture Name (Phase 6).
        """
        start_time = time.time()
        try:
            # 1. Segment text semantically by Heading, Section, or Architecture Name
            default_sec = (additional_metadata or {}).get("section", "Main Content") or "Main Content"
            default_tit = (additional_metadata or {}).get("title", source_file) or source_file
            
            # Helper list of known architectures
            architectures = [
                "Standard RAG", "DeepRAG", "MA-RAG", "Corrective RAG", 
                "Speculative RAG", "Fusion RAG", "RAG-Gym", "Modular RAG", "SAM-RAG"
            ]
            
            lines = text.split('\n')
            segments = []
            current_segment_text = []
            current_section = default_sec
            current_title = default_tit
            
            char_offset = 0
            
            for line in lines:
                line_stripped = line.strip()
                is_header = False
                new_section = None
                
                # Check markdown headings
                h_match = re.match(r'^#{1,6}\s+(.+)$', line_stripped)
                if h_match:
                    is_header = True
                    new_section = h_match.group(1).strip()
                else:
                    # Check numbered headings
                    num_match = re.match(r'^\d+(?:\.\d+)*\s+([A-Z][A-Za-z0-9\s\-\(\)\/\,\.\:\&]+)$', line_stripped)
                    if num_match:
                        is_header = True
                        new_section = num_match.group(1).strip()
                
                # Check architecture names in short lines
                if not is_header and len(line_stripped) < 100:
                    for arch in architectures:
                        if arch.lower() in line_stripped.lower() and (
                            line_stripped.isupper() or 
                            len(line_stripped) == len(arch) or 
                            re.match(r'^(?:\d+[\.\s]+)?' + re.escape(arch) + r'\b', line_stripped, re.IGNORECASE)
                        ):
                            is_header = True
                            new_section = arch
                            break
                            
                if is_header and new_section:
                    if current_segment_text:
                        seg_text = '\n'.join(current_segment_text)
                        segments.append({
                            "text": seg_text,
                            "section": current_section,
                            "title": current_title,
                            "start_char": max(0, char_offset - len(seg_text) - 1)
                        })
                        current_segment_text = []
                    
                    current_section = new_section
                    detected_arch = None
                    for arch in architectures:
                        if arch.lower() in new_section.lower():
                            detected_arch = arch
                            break
                    current_title = detected_arch if detected_arch else new_section
                    
                current_segment_text.append(line)
                char_offset += len(line) + 1
                
            if current_segment_text:
                seg_text = '\n'.join(current_segment_text)
                segments.append({
                    "text": seg_text,
                    "section": current_section,
                    "title": current_title,
                    "start_char": max(0, char_offset - len(seg_text) - 1)
                })

            # 2. Chunk each segment using TokenChunker
            chunks: List[DocumentChunk] = []
            chunk_global_idx = 0
            
            for seg in segments:
                seg_text = seg["text"]
                if not seg_text.strip():
                    continue
                    
                raw_chunks = self.token_chunker.chunk(seg_text)
                
                for chunk_data in raw_chunks:
                    chunk_metadata = additional_metadata.copy() if additional_metadata else {}
                    
                    # Store exact page, section, and title in metadata (Phase 6)
                    chunk_metadata["page"] = page_number
                    chunk_metadata["page_number"] = page_number
                    chunk_metadata["section"] = seg["section"]
                    chunk_metadata["title"] = seg["title"]
                    
                    # Calculate absolute character positions relative to original page text
                    abs_start = seg["start_char"] + chunk_data.start_index
                    abs_end = seg["start_char"] + chunk_data.end_index
                    
                    chunk = DocumentChunk(
                        content=chunk_data.text,
                        source_file=source_file,
                        source_type=source_type,
                        page_number=page_number,
                        chunk_index=chunk_global_idx,
                        start_char=abs_start,
                        end_char=abs_end,
                        metadata=chunk_metadata
                    )
                    chunks.append(chunk)
                    chunk_global_idx += 1

            elapsed = time.time() - start_time
            logger.info(
                f"⚡ Semantic Token chunking (chonkie + heading/arch split): "
                f"{len(chunks)} chunks in {elapsed*1000:.2f}ms"
            )
            return chunks
        except Exception as e:
            logger.error(
                f"❌ TokenChunker failed ({e}). Falling back to standard chunking."
            )
            char_size = self.chunk_size * 4
            char_overlap = self.chunk_overlap * 4
            return self._create_chunks_standard(
                text, source_file, source_type,
                page_number, additional_metadata
            )
    
    def _create_chunks_fast(
        self,
        text: str,
        source_file: str,
        source_type: str,
        page_number: Optional[int] = None,
        additional_metadata: Dict[str, Any] = None
    ) -> List[DocumentChunk]:
        """
        Fast chunking using memchunk (via FastChunker) + post-processing overlap

        Pipeline:
        1️⃣ memchunk delimiter-aware split (no overlap, SIMD-fast)
        2️⃣ Convert to DocumentChunk
        3️⃣ Apply overlap in wrapper (_apply_overlap)
        """
        start_time = time.time()

        try:
            # 1️⃣ Raw delimiter-aware chunks (NO overlap, ultra-fast)
            raw_chunks = self.fast_chunker.chunk(text)

            chunks: List[DocumentChunk] = []

            # 2️⃣ Convert to DocumentChunk objects
            for idx, chunk_data in enumerate(raw_chunks):
                chunk_metadata = additional_metadata.copy() if additional_metadata else {}

                chunk = DocumentChunk(
                    content=chunk_data.text,
                    source_file=source_file,
                    source_type=source_type,
                    page_number=page_number,
                    chunk_index=idx,
                    start_char=chunk_data.start_index,
                    end_char=chunk_data.end_index,
                    metadata=chunk_metadata
                )
                chunks.append(chunk)

            # 3️⃣ Apply overlap (post-chunking, cheap & safe)
            if self.chunk_overlap > 0:
                chunks = self._apply_overlap(chunks)

            elapsed = time.time() - start_time
            speed_mbps = (len(text) / 1024 / 1024) / elapsed if elapsed > 0 else 0

            logger.info(
                f"⚡ Fast chunking (memchunk + overlap): "
                f"{len(chunks)} chunks in {elapsed*1000:.2f}ms "
                f"({speed_mbps:.1f} MB/s, overlap={self.chunk_overlap})"
            )

            return chunks

        except Exception as e:
            logger.error(
                f"❌ FastChunker failed ({e}). Falling back to standard chunking."
            )
            return self._create_chunks_standard(
                text, source_file, source_type,
                page_number, additional_metadata
            )

        
    def _create_chunks_standard(
        self,
        text: str,
        source_file: str,
        source_type: str,
        page_number: Optional[int] = None,
        additional_metadata: Dict[str, Any] = None
    ) -> List[DocumentChunk]:
        """
        Standard chunking (fallback when FastChunker not available)
        
        Uses Python string operations - slower but works without additional dependencies
        """
        chunks = []
        start = 0
        chunk_index = 0
        
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            
            # Try to break at sentence or paragraph boundary
            if end < len(text):
                last_period = text.rfind('.', start, end)
                last_newline = text.rfind('\n', start, end)
                boundary = max(last_period, last_newline)
                
                if boundary > start + self.chunk_size * 0.5:
                    end = boundary + 1
            
            chunk_text = text[start:end].strip()
            
            if chunk_text:
                chunk_metadata = additional_metadata.copy() if additional_metadata else {}
                
                chunk = DocumentChunk(
                    content=chunk_text,
                    source_file=source_file,
                    source_type=source_type,
                    page_number=page_number,
                    chunk_index=chunk_index,
                    start_char=start,
                    end_char=end-1,
                    metadata=chunk_metadata
                )
                chunks.append(chunk)
                chunk_index += 1
            
            start = max(start + self.chunk_size - self.chunk_overlap, end)
            
            if start >= len(text):
                break
        
        return chunks

    def batch_process(self, file_paths: List[str]) -> List[DocumentChunk]:
        """Process multiple files"""
        all_chunks = []
        
        for file_path in file_paths:
            try:
                chunks = self.process_document(file_path)
                all_chunks.extend(chunks)
                logger.info(f"✅ {file_path}: {len(chunks)} chunks")
            except Exception as e:
                logger.error(f"❌ Failed {file_path}: {str(e)}")
                continue
        
        logger.info(
            f"✅ Batch complete: {len(all_chunks)} chunks from {len(file_paths)} files"
        )
        return all_chunks